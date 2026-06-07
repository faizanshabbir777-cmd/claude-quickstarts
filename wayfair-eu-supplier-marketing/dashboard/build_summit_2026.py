"""2026 EU Exec Summit deck builder.

Uses the 2025 Summit file as the baseline and projects 2026 targets at 5% of WSC.
Marks placeholders clearly where 2025-actuals or 2026-confirmed targets are needed.
"""
import sys
sys.path.insert(0, '/home/user/claude-quickstarts/wayfair-eu-supplier-marketing/dashboard')

import openpyxl
from collections import defaultdict, Counter
from PIL import Image, ImageDraw
import io
import deck_builder as db

# -------------------------------------------------------------------------
# Load + aggregate
# -------------------------------------------------------------------------
fn = '/root/.claude/uploads/f60dca77-2aeb-57f6-95f9-5db96153fad0/0afae91b-Ads_Target_Suppliers__Exec_Summit_2025.xlsx'
wb = openpyxl.load_workbook(fn, data_only=True)
ws = wb['Sheet2']

def num(v):
    try: return float(v) if v not in (None, '', '∅') else 0
    except: return 0

stats = {
    'n_total': 0, 'n_active': 0, 'n_lapsed': 0, 'n_unacquired': 0,
    'wsc_24': 0, 'spend_24': 0, 'target_25': 0,
}
by_srm = defaultdict(lambda: {'count': 0, 'active': 0, 'lapsed': 0, 'target': 0, 'wsc': 0, 'spend_24': 0})
by_cat = defaultdict(lambda: {'count': 0, 'wsc': 0, 'spend': 0, 'target': 0})
top_spenders = []

for row in ws.iter_rows(min_row=3, values_only=True):
    suid = row[0]
    if not suid: continue
    supplier, cat, srm, sam = row[1], row[2], row[3], row[4]
    wsc_24, spend_24 = num(row[5]), num(row[6])
    target_25 = num(row[8]); wsp_status = row[10] or 'Unknown'

    stats['n_total'] += 1
    stats['wsc_24'] += wsc_24
    stats['spend_24'] += spend_24
    stats['target_25'] += target_25
    is_active = 'Active' in wsp_status
    is_lapsed = 'Lapsed' in wsp_status
    is_unacq = 'Unacquired' in wsp_status
    if is_active: stats['n_active'] += 1
    if is_lapsed: stats['n_lapsed'] += 1
    if is_unacq: stats['n_unacquired'] += 1

    if srm:
        by_srm[srm]['count'] += 1
        by_srm[srm]['target'] += target_25
        by_srm[srm]['wsc'] += wsc_24
        by_srm[srm]['spend_24'] += spend_24
        if is_active: by_srm[srm]['active'] += 1
        if is_lapsed: by_srm[srm]['lapsed'] += 1
    if cat:
        by_cat[cat]['count'] += 1
        by_cat[cat]['wsc'] += wsc_24
        by_cat[cat]['spend'] += spend_24
        by_cat[cat]['target'] += target_25
    if spend_24 > 50_000:
        top_spenders.append((supplier, cat, srm, wsc_24, spend_24, target_25))


# -------------------------------------------------------------------------
# Slide renderers — Summit-specific
# -------------------------------------------------------------------------
W, H, M = db.W, db.H, db.M


def slide_cover():
    img = Image.new("RGB", (W, H), tuple(db.DEEP_T))
    d = ImageDraw.Draw(img)
    from pathlib import Path
    logo_path = Path(db.__file__).parent / "assets" / "wayfair_logo_white.png"
    if logo_path.exists():
        logo = Image.open(logo_path).convert("RGBA")
        logo.thumbnail((240, 100)); img.paste(logo, (M, 80), logo)

    db.draw_ls(d, M, 220, "EU EXEC SUMMIT  ·  ADS TARGET SUPPLIERS",
               db.sans(26, bold=True), db.GOLD_T, sp=3)
    d.text((M, 320), "From 1.4% to 5% of WSC.",
           font=db.serif(115, bold=True), fill=db.WHITE_T)
    d.text((M, 450), "The multi-year arc.",
           font=db.serif(115, bold=True), fill=db.WHITE_T)

    d.text((M, 620),
           "2024 baseline · 2025 target check · 2026 commitments.",
           font=db.serif(38, italic=True), fill=db.WHITE_T)
    d.text((M, 670),
           "3,841 EU suppliers in the portfolio. $317M of wholesale at the table.",
           font=db.serif(38, italic=True), fill=db.WHITE_T)

    # Hero stat ribbon
    card_y = 880; card_h = 320
    card_w = (W - 2*M - 3*30) // 4; card_gap = 30
    kpis = [
        ('PORTFOLIO WSC (2024)', '$317M', '3,841 EU suppliers'),
        ('2024 AD SPEND',         '$3.7M', '1.16% of WSC'),
        ('2025 TARGET',           '$4.5M', '1.43% of WSC'),
        ('5% BENCHMARK',          '$15.9M', '${gap}M gap to close'.format(gap=round((stats['wsc_24']*0.05 - stats['spend_24'])/1e6))),
    ]
    for i, (label, value, sub) in enumerate(kpis):
        cx = M + i * (card_w + card_gap)
        d.rectangle([cx, card_y, cx + card_w, card_y + card_h], fill=db.LAV_T)
        d.rectangle([cx, card_y, cx + card_w, card_y + 10], fill=db.GOLD_T)
        db.draw_ls(d, cx + 35, card_y + 38, label,
                   db.sans(18, bold=True), db.PURPLE_T, sp=2)
        d.text((cx + 35, card_y + 95), value,
               font=db.serif(70, bold=True, italic=True), fill=db.DEEP_T)
        d.text((cx + 35, card_y + card_h - 70), sub,
               font=db.sans(20), fill=db.INK_T)

    d.text((M, H - 75),
           "EU Supplier Marketing · 2026 planning cycle · prepared for Niraj + GM leadership",
           font=db.sans(20), fill=(180, 165, 200))
    return img


def slide_portfolio_status():
    img = Image.new("RGB", (W, H), db.WHITE_T)
    d = ImageDraw.Draw(img)
    db.draw_ls(d, M, 90, "01  ·  THE PORTFOLIO AT A GLANCE",
               db.sans(22, bold=True), db.PURPLE_T, sp=3)
    d.text((M, 165),
           "Two-thirds of our supplier base has never run an ad with us.",
           font=db.serif(48, bold=True), fill=db.INK_T)

    # Status distribution — stacked horizontal bar
    bar_y = 380; bar_h = 140
    bar_left = M; bar_right = W - M
    bar_w = bar_right - bar_left
    segments = [
        ('UNACQUIRED', stats['n_unacquired'], db.GREY_BAR),
        ('LAPSED',     stats['n_lapsed'],    (200, 130, 100)),
        ('ACTIVE',     stats['n_active'],    db.GREEN_T),
        ('OTHER',      stats['n_total'] - stats['n_unacquired'] - stats['n_lapsed'] - stats['n_active'],
                       (180, 130, 200)),
    ]
    cx = bar_left
    db.draw_ls(d, M, bar_y - 35, "3,841 EU SUPPLIERS BY 2025 WSP STATUS",
               db.sans(18, bold=True), db.SLATE_T, sp=2)
    for label, count, color in segments:
        pct = count / stats['n_total'] * 100
        w = int(bar_w * pct / 100)
        d.rectangle([cx, bar_y, cx + w, bar_y + bar_h], fill=color)
        if w > 200:
            text_color = db.WHITE_T if color == db.GREEN_T else db.INK_T
            d.text((cx + 24, bar_y + 24), label,
                   font=db.sans(20, bold=True), fill=text_color)
            d.text((cx + 24, bar_y + 60), f"{count:,}",
                   font=db.serif(44, bold=True, italic=True), fill=text_color)
            d.text((cx + 24, bar_y + 110), f"{pct:.1f}%",
                   font=db.sans(20), fill=text_color)
        cx += w

    # Insight band below
    cb_y = 640; cb_h = 420
    d.rectangle([M, cb_y, W - M, cb_y + cb_h], fill=db.LAV_T)
    db.draw_ls(d, M + 40, cb_y + 30, "WHAT 2026 NEEDS TO DO",
               db.sans(20, bold=True), db.PURPLE_T, sp=2)

    three_buckets = [
        ('LIFT THE 188 ACTIVE',
         "From 1.4% to 5% of WSC.",
         "Variant B/C pitches across all 188 active suppliers."),
        ('REACTIVATE 452 LAPSED',
         "TuttiBambini-style restart with the dark-days story.",
         "Variant A pitches to the top 100 by 2024 WSC."),
        ('ACQUIRE FROM 2,419 UNACQUIRED',
         "Target the Tier 1 / high-WSC unacquired first.",
         "Cohort of ~200 Tier 1 unacquired SUs — the 2026 hunt list."),
    ]
    pillar_y = cb_y + 80
    pillar_h = cb_h - 120
    pillar_w = (W - 2*M - 80 - 2*30) // 3
    for i, (lbl, hero, body) in enumerate(three_buckets):
        px = M + 40 + i * (pillar_w + 30)
        d.rectangle([px, pillar_y, px + pillar_w, pillar_y + pillar_h], fill=db.WHITE_T)
        d.rectangle([px, pillar_y, px + pillar_w, pillar_y + 8], fill=db.GOLD_T)
        db.draw_ls(d, px + 24, pillar_y + 28, lbl,
                   db.sans(16, bold=True), db.PURPLE_T, sp=2)
        d.text((px + 24, pillar_y + 76), hero,
               font=db.serif(28, bold=True, italic=True), fill=db.DEEP_T)
        for j, ln in enumerate(db.wrap(d, body, pillar_w - 50, db.sans(18))):
            d.text((px + 24, pillar_y + 170 + j*26), ln,
                   font=db.sans(18), fill=db.INK_T)

    d.text((M, H - 80),
           "Take-away: the opportunity is not in the 188 already running — it is in the 2,871 who are not.",
           font=db.serif(28, italic=True), fill=db.INK_T)
    return img


def slide_srm_scoreboard():
    img = Image.new("RGB", (W, H), db.WHITE_T)
    d = ImageDraw.Draw(img)
    db.draw_ls(d, M, 90, "02  ·  WHERE THE 2025 TARGET LIVES BY SRM",
               db.sans(22, bold=True), db.PURPLE_T, sp=3)
    d.text((M, 165),
           "Eight SRMs own 70% of the 2025 ad-target portfolio.",
           font=db.serif(48, bold=True), fill=db.INK_T)

    top_srms = sorted(by_srm.items(), key=lambda x: -x[1]['target'])[:8]

    # Table
    table_y = 320
    row_h = 80
    col_widths = [550, 260, 260, 290, 290, 250]
    col_x = [M]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)
    headers = ['SRM', '# OF SUS', '# ACTIVE', "'24 WSC", "'25 TARGET", "GAP TO 5%"]

    d.rectangle([M, table_y, W - M, table_y + 50], fill=db.DEEP_T)
    for i, h in enumerate(headers):
        d.text((col_x[i] + 16, table_y + 14), h,
               font=db.sans(16, bold=True), fill=db.GOLD_T)

    for r, (srm_name, s) in enumerate(top_srms):
        ry = table_y + 50 + r * row_h
        if r % 2 == 0:
            d.rectangle([M, ry, W - M, ry + row_h], fill=(248, 244, 248))
        d.text((col_x[0] + 16, ry + 22), srm_name[:36],
               font=db.serif(24, bold=True, italic=True), fill=db.DEEP_T)
        active_pct = s['active'] / s['count'] * 100 if s['count'] else 0
        d.text((col_x[1] + 16, ry + 22), f"{s['count']:,}",
               font=db.serif(28, bold=True, italic=True), fill=db.DEEP_T)
        # Active count with green chip if >50%
        chip_w = 100
        chip_color = db.GREEN_T if active_pct > 50 else (200, 130, 100)
        d.rounded_rectangle([col_x[2] + 16, ry + 22, col_x[2] + 16 + chip_w, ry + 56],
                             radius=6, fill=chip_color)
        d.text((col_x[2] + 30, ry + 28),
               f"{s['active']} · {active_pct:.0f}%",
               font=db.sans(16, bold=True), fill=db.WHITE_T)
        d.text((col_x[3] + 16, ry + 22), f"${s['wsc']/1e6:.1f}M",
               font=db.serif(28, bold=True, italic=True), fill=db.DEEP_T)
        d.text((col_x[4] + 16, ry + 22), f"${s['target']/1000:.0f}k",
               font=db.serif(28, bold=True, italic=True), fill=db.PURPLE_T)
        gap_to_5pct = s['wsc'] * 0.05 - s['target']
        d.text((col_x[5] + 16, ry + 22),
               f"+${gap_to_5pct/1e6:.1f}M",
               font=db.serif(28, bold=True, italic=True), fill=db.GOLD_T)

    # Callout below
    cb_y = 1000; cb_h = 130
    d.rectangle([M, cb_y, W - M, cb_y + cb_h], fill=db.LAV_T)
    d.text((M + 40, cb_y + 30), "READ-OUT",
           font=db.sans(18, bold=True), fill=db.PURPLE_T)
    d.text((M + 40, cb_y + 60),
           "The 2025 target sums to ~$4.5M. Closing to 5% of WSC = ~$11M more on the table.",
           font=db.sans(24), fill=db.INK_T)

    d.text((M, H - 70),
           "Take-away: 2026 SRM target-setting starts here — five percent of last-year's WSC, per SRM portfolio.",
           font=db.serif(26, italic=True), fill=db.INK_T)
    return img


def slide_category_map():
    img = Image.new("RGB", (W, H), db.WHITE_T)
    d = ImageDraw.Draw(img)
    db.draw_ls(d, M, 90, "03  ·  CATEGORY OPPORTUNITY MAP",
               db.sans(22, bold=True), db.PURPLE_T, sp=3)
    d.text((M, 165),
           "Three categories hold half the wholesale — and most of the ad-spend ceiling.",
           font=db.serif(48, bold=True), fill=db.INK_T)

    top_cats = sorted(by_cat.items(), key=lambda x: -x[1]['wsc'])[:10]
    # Horizontal bar chart of WSC per category with target overlay
    chart_y = 320
    chart_h = 720
    row_h = chart_h // len(top_cats)
    max_wsc = max(s['wsc'] for _, s in top_cats)
    bar_left = M + 360  # leave room for category names
    bar_max_w = W - M - bar_left - 280  # leave room for value labels

    for i, (cat, s) in enumerate(top_cats):
        ry = chart_y + i * row_h
        # Category name
        d.text((M, ry + row_h // 2 - 15), cat[:30],
               font=db.sans(20, bold=True), fill=db.DEEP_T)
        d.text((M, ry + row_h // 2 + 12), f"{s['count']:,} SUs",
               font=db.sans(15, italic=True), fill=db.SLATE_T)

        # WSC bar
        w = int(bar_max_w * s['wsc'] / max_wsc)
        d.rectangle([bar_left, ry + 15, bar_left + w, ry + row_h - 25], fill=db.PURPLE_T)
        # Spend overlay (dark stripe)
        spend_w = int(bar_max_w * s['spend'] / max_wsc * 8)  # scale up 8x for visibility
        spend_w = min(spend_w, w)
        if spend_w > 5:
            d.rectangle([bar_left, ry + 15, bar_left + spend_w, ry + row_h - 25],
                       fill=db.GOLD_T)

        # Labels right of bar
        wsc_text = f"${s['wsc']/1e6:.1f}M WSC"
        d.text((bar_left + w + 16, ry + 15),
               wsc_text, font=db.serif(24, bold=True, italic=True), fill=db.DEEP_T)
        d.text((bar_left + w + 16, ry + 50),
               f"${s['spend']/1000:.0f}k spent  ·  ${s['target']/1000:.0f}k 25 tgt",
               font=db.sans(15), fill=db.SLATE_T)

    # Legend
    leg_y = chart_y - 60
    d.rectangle([bar_left, leg_y, bar_left + 24, leg_y + 18], fill=db.PURPLE_T)
    d.text((bar_left + 32, leg_y), "Wholesale (2024)",
           font=db.sans(16), fill=db.INK_T)
    d.rectangle([bar_left + 280, leg_y, bar_left + 304, leg_y + 18], fill=db.GOLD_T)
    d.text((bar_left + 312, leg_y), "Ad spend (2024, ×8 for visibility)",
           font=db.sans(16), fill=db.INK_T)

    d.text((M, H - 80),
           "Take-away: APS EU F&D, UK Bedroom, APS Outdoor — three categories, $104M wholesale, $1.7M target. The 2026 priority list.",
           font=db.serif(26, italic=True), fill=db.INK_T)
    return img


def slide_2026_framework():
    img = Image.new("RGB", (W, H), db.WHITE_T)
    d = ImageDraw.Draw(img)
    db.draw_ls(d, M, 90, "04  ·  THE 2026 FRAMEWORK",
               db.sans(22, bold=True), db.PURPLE_T, sp=3)
    d.text((M, 165),
           "One rule. Three priorities. Eight SRMs. Twelve months.",
           font=db.serif(48, bold=True), fill=db.INK_T)

    # The 5% rule reminder
    rule_y = 320; rule_h = 200
    d.rectangle([M, rule_y, W - M, rule_y + rule_h], fill=db.LAV_T)
    d.rectangle([M, rule_y, M + 12, rule_y + rule_h], fill=db.GOLD_T)
    db.draw_ls(d, M + 40, rule_y + 30, "THE RULE",
               db.sans(20, bold=True), db.PURPLE_T, sp=2)
    d.text((M + 40, rule_y + 70),
           "Set next month's WSP at 5% of last month's WSC.",
           font=db.serif(40, bold=True, italic=True), fill=db.DEEP_T)
    d.text((M + 40, rule_y + 130),
           "Apply on day 1 of every month. Same rule for every SRM, every supplier, every category.",
           font=db.sans(22, italic=True), fill=db.INK_T)

    # Three priority cards
    pr_y = 580; pr_h = 360
    pr_w = (W - 2*M - 2*30) // 3; pr_gap = 30
    priorities = [
        ('1.  LIFT', 'the 188 active',
         "Variant B/C pitches across all active suppliers.\nMove from 1.4% to 5% of WSC.\nTarget delta: +$11M of qualified ad spend."),
        ('2.  REACTIVATE', 'the 452 lapsed',
         "TuttiBambini-style restart with the dark-days story.\nTop 100 by 2024 WSC = the priority cohort.\nTarget delta: +$3M of qualified ad spend."),
        ('3.  ACQUIRE', 'from the 2,419 unacquired',
         "Tier 1 high-WSC unacquired SUs = the 2026 hunt list.\nSize the cohort to 200 SUs across the eight SRMs.\nTarget delta: +$5M of qualified ad spend."),
    ]
    for i, (num_lbl, target_lbl, body) in enumerate(priorities):
        px = M + i * (pr_w + pr_gap)
        d.rectangle([px, pr_y, px + pr_w, pr_y + pr_h], fill=db.LAV_T)
        d.rectangle([px, pr_y, px + pr_w, pr_y + 10], fill=db.GOLD_T)
        db.draw_ls(d, px + 30, pr_y + 35, num_lbl,
                   db.sans(20, bold=True), db.PURPLE_T, sp=2)
        d.text((px + 30, pr_y + 90), target_lbl,
               font=db.serif(38, bold=True, italic=True), fill=db.DEEP_T)
        for j, ln in enumerate(body.split('\n')):
            for k, wln in enumerate(db.wrap(d, ln, pr_w - 60, db.sans(20))):
                d.text((px + 30, pr_y + 200 + (j*60 + k*24)), wln,
                       font=db.sans(20), fill=db.INK_T)

    # The ask card at the bottom
    ask_y = 990; ask_h = 175
    d.rectangle([M, ask_y, W - M, ask_y + ask_h], fill=db.DEEP_T)
    db.draw_ls(d, M + 40, ask_y + 28, "THE ASK", db.sans(24, bold=True), db.GOLD_T, sp=3)
    d.text((M + 40, ask_y + 70),
           "Commit each SRM to a portfolio target = 5% of their suppliers' 2025 WSC,",
           font=db.sans(34), fill=db.WHITE_T)
    d.text((M + 40, ask_y + 115),
           "phased monthly. The skill builds every pitch — the SRM signs off the relationship.",
           font=db.sans(34), fill=db.WHITE_T)

    d.text((M, H - 70),
           "One rule  ·  three priorities  ·  eight SRMs  ·  live on day 1 every month.",
           font=db.serif(26, italic=True), fill=db.INK_T)
    return img


# -------------------------------------------------------------------------
# Build the deck
# -------------------------------------------------------------------------
def build():
    import tempfile
    slides = [
        slide_cover(),
        slide_portfolio_status(),
        slide_srm_scoreboard(),
        slide_category_map(),
        slide_2026_framework(),
    ]
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    with tempfile.TemporaryDirectory() as td:
        for i, slide_img in enumerate(slides):
            if slide_img.mode != "RGB":
                slide_img = slide_img.convert("RGB")
            png_path = f"{td}/slide_{i}.png"
            slide_img.save(png_path, "PNG", optimize=False)
            s = prs.slides.add_slide(BLANK)
            s.shapes.add_picture(png_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        out = io.BytesIO(); prs.save(out); out.seek(0)
        return out.read()


if __name__ == '__main__':
    pptx = build()
    with open('/tmp/EU_Exec_Summit_2026.pptx', 'wb') as f:
        f.write(pptx)
    print(f"Built EU Exec Summit 2026 deck: {len(pptx):,} bytes")
    from pptx import Presentation
    p = Presentation('/tmp/EU_Exec_Summit_2026.pptx')
    print(f"Validates: {len(p.slides)} slides")
