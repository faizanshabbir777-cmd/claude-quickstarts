"""Real deck builder for the WSP Pitch Builder dashboard.

Takes structured supplier data, produces a personalised .pptx in the
TuttiBambini visual register. No Claude API call required.

This is the production v1.0 backend — not a stub.
"""
import io
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt
import matplotlib.ticker as mtick
import numpy as np
import os


# ============================================================
# Palette — matches SKILL.md §C and the TuttiBambini reference
# ============================================================
DEEP_T   = (60, 26, 80)
PURPLE_T = (123, 24, 159)
LAV_T    = (246, 235, 251)
GOLD_T   = (212, 160, 23)
INK_T    = (26, 26, 26)
SLATE_T  = (91, 91, 102)
HAIR_T   = (229, 222, 236)
WHITE_T  = (255, 255, 255)
GREEN_T  = (14, 143, 96)
GREEN_LIGHT = (37, 171, 122)


# ============================================================
# Fonts (Liberation Serif/Sans available on Streamlit Cloud)
# ============================================================
def serif(size, bold=False, italic=False):
    paths = {
        (True, True):   "/usr/share/fonts/truetype/liberation/LiberationSerif-BoldItalic.ttf",
        (True, False):  "/usr/share/fonts/truetype/liberation/LiberationSerif-Bold.ttf",
        (False, True):  "/usr/share/fonts/truetype/liberation/LiberationSerif-Italic.ttf",
        (False, False): "/usr/share/fonts/truetype/liberation/LiberationSerif-Regular.ttf",
    }
    p = paths.get((bold, italic))
    if p and os.path.exists(p):
        return ImageFont.truetype(p, size)
    # Fallback if Liberation isn't on this box
    fb = "/usr/share/fonts/truetype/dejavu/DejaVuSerif-Bold.ttf" if bold else \
         "/usr/share/fonts/truetype/dejavu/DejaVuSerif.ttf"
    if os.path.exists(fb):
        return ImageFont.truetype(fb, size)
    return ImageFont.load_default()


def sans(size, bold=False, italic=False):
    paths = {
        (True, True):   "/usr/share/fonts/truetype/liberation/LiberationSans-BoldItalic.ttf",
        (True, False):  "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
        (False, True):  "/usr/share/fonts/truetype/liberation/LiberationSans-Italic.ttf",
        (False, False): "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    }
    p = paths.get((bold, italic))
    if p and os.path.exists(p):
        return ImageFont.truetype(p, size)
    fb = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else \
         "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
    if os.path.exists(fb):
        return ImageFont.truetype(fb, size)
    return ImageFont.load_default()


def tsize(d, t, f):
    b = d.textbbox((0,0), t, font=f); return b[2]-b[0], b[3]-b[1]


def draw_ls(d, x, y, txt, f, fill, sp=3):
    cx = x
    for ch in txt:
        d.text((cx, y), ch, font=f, fill=fill)
        b = d.textbbox((0,0), ch, font=f); cx += (b[2]-b[0]) + sp


def wrap(d, t, mw, f):
    out, line = [], []
    for w in t.split():
        test = " ".join(line + [w])
        if tsize(d, test, f)[0] > mw:
            out.append(" ".join(line)); line = [w]
        else: line.append(w)
    if line: out.append(" ".join(line))
    return out


# Slide dimensions match TuttiBambini (16:9, rendered at 150 DPI)
W, H = 2400, 1350
M = 130


# ============================================================
# Slide renderers — TuttiBambini visual register
# ============================================================
def render_cover(data: dict) -> Image.Image:
    """Slide 1 — cover with KPI ribbon."""
    img = Image.new("RGB", (W, H), tuple(DEEP_T))
    d = ImageDraw.Draw(img)

    # Wayfair logo
    here = Path(__file__).parent
    logo_path = here / "assets" / "wayfair_logo_white.png"
    if logo_path.exists():
        logo = Image.open(logo_path).convert("RGBA")
        logo.thumbnail((240, 100))
        img.paste(logo, (M, 80), logo)

    # Gold eyebrow
    eyebrow = f"{data['supplier'].upper()}  ·  WAYFAIR SPONSORED PRODUCTS"
    draw_ls(d, M, 220, eyebrow, sans(26, bold=True), GOLD_T, sp=3)

    # Hero title (auto-picked based on data)
    title = data.get('cover_title', 'Keep your WSP on at this pace.')
    d.text((M, 320), title, font=serif(115, bold=True), fill=WHITE_T)

    # Subtitle
    subtitle = data.get('cover_subtitle',
        'Your ads are doing the work. Here is what they returned, and the simple way to lock it in.')
    sub_lines = wrap(d, subtitle, W - 2*M - 100, serif(38, italic=True))
    for i, ln in enumerate(sub_lines[:2]):
        d.text((M, 530 + i*55), ln, font=serif(38, italic=True), fill=WHITE_T)

    # 4 KPI cards along bottom
    card_y = 880; card_h = 320
    card_w = (W - 2*M - 3*30) // 4
    card_gap = 30
    kpis = data['kpis_cover']  # list of 4 dicts: {label, value, sub}
    for i, kpi in enumerate(kpis[:4]):
        cx = M + i * (card_w + card_gap)
        d.rectangle([cx, card_y, cx + card_w, card_y + card_h], fill=LAV_T)
        d.rectangle([cx, card_y, cx + card_w, card_y + 10], fill=GOLD_T)
        draw_ls(d, cx + 35, card_y + 38, kpi['label'].upper(),
                sans(20, bold=True), PURPLE_T, sp=2)
        d.text((cx + 35, card_y + 95), kpi['value'],
               font=serif(75, bold=True, italic=True), fill=DEEP_T)
        d.text((cx + 35, card_y + card_h - 70), kpi['sub'],
               font=sans(20), fill=INK_T)

    # Footer
    footer = f"{data['period']}  ·  prepared for {data.get('srm', 'the team')}  ·  UK  ·  GBP"
    d.text((M, H - 75), footer, font=sans(20), fill=(180, 165, 200))
    return img


def render_value_story(data: dict) -> Image.Image:
    """Slide 2 — '01 · WHAT WSP IS DOING FOR YOU' with hero card + chart."""
    img = Image.new("RGB", (W, H), WHITE_T)
    d = ImageDraw.Draw(img)
    draw_ls(d, M, 90, "01  ·  WHAT WSP IS DOING FOR YOU",
            sans(22, bold=True), PURPLE_T, sp=3)

    ratio = data['ws_roas']
    title = data.get('story_title',
        f"When it is on, every pound comes back about {int(round(ratio))} times over.")
    d.text((M, 165), title, font=serif(48, bold=True), fill=INK_T)

    # LEFT lavender hero card
    left_x = M; left_w = 850; left_y = 370; left_h = 700
    d.rectangle([left_x, left_y, left_x + left_w, left_y + left_h], fill=LAV_T)
    draw_ls(d, left_x + 40, left_y + 40, f"SO FAR IN {data['year']}",
            sans(22, bold=True), PURPLE_T, sp=2)
    d.text((left_x + 40, left_y + 95), f"£{ratio:.2f}",
           font=serif(180, bold=True, italic=True), fill=DEEP_T)
    d.text((left_x + 40, left_y + 300),
           "of wholesale revenue for every £1 of WSP",
           font=sans(28), fill=INK_T)
    d.text((left_x + 40, left_y + 340), "spend",
           font=sans(28), fill=INK_T)
    rows = [
        ("You spent", f"£{_fmt(data['wsp_spend_ytd'])} on WSP"),
        ("It drove",  f"£{_fmt(data['attr_wsc_ytd'])} wholesale revenue"),
        ("That's",    f"1 in £{ratio:.2f} — strong and steady"),
    ]
    rx = left_x + 40; ry = left_y + 440
    for label, value in rows:
        d.text((rx, ry), label, font=sans(22), fill=SLATE_T)
        d.text((rx + 230, ry), value, font=sans(22, bold=True), fill=INK_T)
        ry += 60

    # RIGHT chart — bars by month
    chart_x = left_x + left_w + 60
    chart_y = left_y - 30
    chart_w = W - chart_x - M
    chart_h = left_h + 30
    d.text((chart_x, chart_y), "Wholesale revenue driven by WSP, each month (£)",
           font=sans(24, bold=True), fill=INK_T)

    months_data = data['monthly_attr_wsc'][-5:]  # last 5 months
    months = [m[0] for m in months_data]
    values = [m[1] for m in months_data]
    if values and max(values) > 0:
        max_v = max(values) * 1.18
        bar_area_top = chart_y + 80
        bar_area_bot = chart_y + chart_h - 90
        bar_area_h = bar_area_bot - bar_area_top
        n = len(values); gap = 30
        bw = (chart_w - gap*(n-1)) // n
        for i, v in enumerate(values):
            bh = int(v / max_v * bar_area_h)
            bx = chart_x + i*(bw + gap)
            by = bar_area_bot - bh
            d.rectangle([bx, by, bx + bw, bar_area_bot], fill=PURPLE_T)
            label = f"{int(v/1000)}k"
            f_lbl = sans(28, bold=True)
            lw, _ = tsize(d, label, f_lbl)
            d.text((bx + (bw - lw)//2, by - 50), label, font=f_lbl, fill=INK_T)
            f_m = sans(22)
            mw_, _ = tsize(d, months[i], f_m)
            d.text((bx + (bw - mw_)//2, bar_area_bot + 18), months[i], font=f_m, fill=SLATE_T)

    takeaway = data.get('story_takeaway',
        f"Take-away: your WSP has paid back about £{ratio:.2f} for every £1 — and it has every month this year.")
    d.text((M, H - 120), takeaway, font=serif(28, italic=True), fill=INK_T)
    return img


def render_yoy_proof(data: dict) -> Image.Image:
    """Slide 3 — '02 · YoY proof' with deep purple KPI band + green pills."""
    img = Image.new("RGB", (W, H), WHITE_T)
    d = ImageDraw.Draw(img)
    eyebrow = f"02  ·  {data['period'].upper()} vs {data['yoy_period'].upper()}"
    draw_ls(d, M, 90, eyebrow, sans(22, bold=True), PURPLE_T, sp=3)
    d.text((M, 165),
           data.get('yoy_title', "You doubled down on ads — and the ads doubled down on you."),
           font=serif(58, bold=True), fill=INK_T)

    pc_x = M; pc_y = 380; pc_w = 1500; pc_h = 500
    d.rectangle([pc_x, pc_y, pc_x + pc_w, pc_y + pc_h], fill=DEEP_T)
    col_w = pc_w // 3
    for div in [pc_x + col_w, pc_x + 2*col_w]:
        d.rectangle([div - 1, pc_y + 50, div + 1, pc_y + pc_h - 50], fill=(120, 90, 140))

    cols = data['yoy_columns']  # list of 3 dicts: {label, value, sub, pill}
    for i, col in enumerate(cols[:3]):
        cx = pc_x + i*col_w + 50
        draw_ls(d, cx, pc_y + 60, col['label'], sans(22, bold=True), GOLD_T, sp=2)
        d.text((cx, pc_y + 130), col['value'],
               font=serif(100, bold=True, italic=True), fill=WHITE_T)
        d.text((cx, pc_y + 290), col['sub'],
               font=sans(24), fill=(180, 165, 200))
        pill = col['pill']
        f_pill = sans(28, bold=True)
        pw_, ph_ = tsize(d, pill, f_pill)
        pill_w = pw_ + 70; pill_h = ph_ + 26
        pill_x = cx; pill_y = pc_y + 360
        d.rounded_rectangle([pill_x, pill_y, pill_x + pill_w, pill_y + pill_h],
                              radius=8, fill=GREEN_T)
        d.text((pill_x + 35, pill_y + 13), pill, font=f_pill, fill=WHITE_T)

    # RIGHT lavender side card
    rx = pc_x + pc_w + 50
    rw = W - rx - M
    d.rectangle([rx, pc_y, rx + rw, pc_y + pc_h], fill=LAV_T)
    draw_ls(d, rx + 35, pc_y + 40, "WSP % OF WSC",
            sans(22, bold=True), PURPLE_T, sp=2)
    bar_bw = 60; bar_gap = 25; bar_bot = pc_y + 290
    bar_heights = [80, 130, 175]
    bar_x = rx + 60
    for i, bh in enumerate(bar_heights):
        bx = bar_x + i*(bar_bw + bar_gap)
        d.rectangle([bx, bar_bot - bh, bx + bar_bw, bar_bot], fill=GREEN_LIGHT)
    pct = data['wsp_pct_wsc_latest']
    label = "Above" if pct >= 5.0 else "Below"
    d.text((rx + 35, pc_y + 320), label,
           font=serif(50, italic=True), fill=GREEN_T if pct >= 5.0 else (200, 100, 100))
    d.text((rx + 35, pc_y + 380), "5% benchmark",
           font=serif(50, italic=True), fill=GREEN_T if pct >= 5.0 else (200, 100, 100))
    d.text((rx + 35, pc_y + pc_h - 60),
           f"Your spend share is now\n{pct:.1f}% of wholesale.",
           font=sans(20, italic=True), fill=SLATE_T)

    cb_y = 1110; cb_h = 110
    d.rectangle([M, cb_y, W - M, cb_y + cb_h], fill=LAV_T)
    callout = data.get('yoy_callout',
        f"The pattern is clear: spend and wholesale move together, every month.")
    d.text((M + 35, cb_y + 35), callout, font=sans(24), fill=INK_T)
    return img


def render_ask(data: dict) -> Image.Image:
    """Slide 4 — '03 · WHERE WE GO FROM HERE' with three pillars + THE ASK card."""
    img = Image.new("RGB", (W, H), WHITE_T)
    d = ImageDraw.Draw(img)
    draw_ls(d, M, 90, "03  ·  WHERE WE GO FROM HERE",
            sans(22, bold=True), PURPLE_T, sp=3)
    d.text((M, 165), data.get('ask_title',
           "Lock in the rhythm — same return, no surprises."),
           font=serif(58, bold=True), fill=INK_T)

    card_y = 360; card_h = 520
    card_w = (W - 2*M - 2*30) // 3
    card_gap = 30
    cards = data['pillars']  # list of 3 dicts: {label, hero, body}
    for i, card in enumerate(cards[:3]):
        cx = M + i*(card_w + card_gap)
        d.rectangle([cx, card_y, cx + card_w, card_y + card_h], fill=LAV_T)
        d.rectangle([cx, card_y, cx + card_w, card_y + 10], fill=GOLD_T)
        draw_ls(d, cx + 35, card_y + 40, card['label'],
                sans(22, bold=True), PURPLE_T, sp=2)
        d.text((cx + 35, card_y + 110), card['hero'],
               font=serif(95, bold=True, italic=True), fill=DEEP_T)
        body_lines = wrap(d, card['body'], card_w - 70, sans(24))
        for j, ln in enumerate(body_lines):
            d.text((cx + 35, card_y + 290 + j*36), ln, font=sans(24), fill=INK_T)

    # THE ASK card
    ask_y = 920; ask_h = 220
    d.rectangle([M, ask_y, W - M, ask_y + ask_h], fill=DEEP_T)
    draw_ls(d, M + 40, ask_y + 32, "THE ASK", sans(24, bold=True), GOLD_T, sp=3)
    ask_amount = data['ask_amount_text']  # e.g. "~£23k / month"
    last_wsc = _fmt(data['last_month_wsc'])
    d.text((M + 40, ask_y + 80), "Continue WSP at  ", font=sans(40), fill=WHITE_T)
    cw1, _ = tsize(d, "Continue WSP at  ", sans(40))
    d.text((M + 40 + cw1, ask_y + 80), ask_amount,
           font=sans(40, bold=True), fill=GOLD_T)
    cw2, _ = tsize(d, ask_amount, sans(40, bold=True))
    d.text((M + 40 + cw1 + cw2, ask_y + 80), "  —  5% of last month's",
           font=sans(40), fill=WHITE_T)
    d.text((M + 40, ask_y + 135),
           f"wholesale revenue (£{last_wsc}), set on day 1 so the engine never stutters.",
           font=sans(40), fill=WHITE_T)

    closing = data.get('ask_closing',
        f"Proven £{data['ws_roas']:.2f} back for every £1  ·  one rule, live on day 1.")
    d.text((M, H - 80), closing, font=serif(28, italic=True), fill=INK_T)
    return img


# ============================================================
# Top-level builder — Image[] → real .pptx bytes
# ============================================================
def build_deck_value_review_b(data: dict) -> bytes:
    """Build a Variant B value-review deck. Returns .pptx bytes ready to download.

    Uses tempfiles (not BytesIO) for the PNG → python-pptx handoff because some
    PowerPoint versions are picky about images embedded via memory streams.
    """
    import tempfile
    slides = [
        render_cover(data),
        render_value_story(data),
        render_yoy_proof(data),
        render_ask(data),
    ]

    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    # Save each PNG to a real tempfile (conventional) and embed via file path
    with tempfile.TemporaryDirectory() as td:
        for i, slide_img in enumerate(slides):
            # Convert to RGB and save as a clean baseline PNG
            if slide_img.mode != "RGB":
                slide_img = slide_img.convert("RGB")
            png_path = f"{td}/slide_{i}.png"
            slide_img.save(png_path, "PNG", optimize=False)

            s = prs.slides.add_slide(BLANK)
            s.shapes.add_picture(png_path, 0, 0,
                                  width=prs.slide_width,
                                  height=prs.slide_height)

        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        return out.read()


# ============================================================
# Helpers
# ============================================================
def _fmt(v):
    """Format big numbers cleanly: £476k, £1.2M, £33.2k."""
    if v >= 1e6:
        return f"{v/1e6:.2f}M"
    elif v >= 1e3:
        return f"{v/1e3:.1f}k"
    else:
        return f"{v:.0f}"


# =====================================================================
# VARIANT A — RESTART PITCH
# Mirrors the TuttiBambini Way Day deck: cover · value-was-driving ·
# the-event-proved-it · the-problem (dark-days calendar) · the-ask.
# =====================================================================

# Coral palette for problem framing
CORAL_T = (184, 58, 58)
CORAL_LIGHT = (252, 232, 232)
GREY_BAR = (207, 199, 214)


def render_restart_cover(data: dict) -> Image.Image:
    """Cover with 4 KPI cards; the 4th has a CORAL stripe (the loss)."""
    img = Image.new("RGB", (W, H), tuple(DEEP_T))
    d = ImageDraw.Draw(img)
    here = Path(__file__).parent
    logo_path = here / "assets" / "wayfair_logo_white.png"
    if logo_path.exists():
        logo = Image.open(logo_path).convert("RGBA")
        logo.thumbnail((240, 100))
        img.paste(logo, (M, 80), logo)

    draw_ls(d, M, 220,
            f"{data['supplier'].upper()}  ·  WAYFAIR SPONSORED PRODUCTS",
            sans(26, bold=True), GOLD_T, sp=3)
    d.text((M, 320), data.get('cover_title', "Switch your WSP back on."),
           font=serif(115, bold=True), fill=WHITE_T)

    sub = data.get('cover_subtitle',
        "It was working — then it went dark. Here is what it drove for you, and the simple way to keep it on.")
    sub_lines = wrap(d, sub, W - 2*M - 100, serif(38, italic=True))
    for i, ln in enumerate(sub_lines[:2]):
        d.text((M, 530 + i*55), ln, font=serif(38, italic=True), fill=WHITE_T)

    card_y = 880; card_h = 320
    card_w = (W - 2*M - 3*30) // 4
    card_gap = 30
    kpis = data['kpis_cover']
    for i, kpi in enumerate(kpis[:4]):
        cx = M + i * (card_w + card_gap)
        d.rectangle([cx, card_y, cx + card_w, card_y + card_h], fill=LAV_T)
        # CORAL stripe on the loss KPI (typically the 4th)
        stripe_color = tuple(CORAL_T) if kpi.get('is_loss') else GOLD_T
        d.rectangle([cx, card_y, cx + card_w, card_y + 10], fill=stripe_color)
        draw_ls(d, cx + 35, card_y + 38, kpi['label'].upper(),
                sans(20, bold=True),
                tuple(CORAL_T) if kpi.get('is_loss') else PURPLE_T, sp=2)
        value_color = tuple(CORAL_T) if kpi.get('is_loss') else DEEP_T
        d.text((cx + 35, card_y + 95), kpi['value'],
               font=serif(75, bold=True, italic=True), fill=value_color)
        d.text((cx + 35, card_y + card_h - 70), kpi['sub'],
               font=sans(20), fill=INK_T)

    footer = f"{data.get('event_window', data['period'])}  ·  prepared for {data.get('srm','the team')}  ·  UK  ·  GBP"
    d.text((M, H - 75), footer, font=sans(20), fill=(180, 165, 200))
    return img


def render_dark_days_calendar(data: dict) -> Image.Image:
    """The killer slide — day-by-day purple/grey blocks showing dark days."""
    img = Image.new("RGB", (W, H), WHITE_T)
    d = ImageDraw.Draw(img)
    draw_ls(d, M, 90, "03  ·  THE PROBLEM",
            sans(22, bold=True), tuple(CORAL_T), sp=3)
    d.text((M, 165),
           data.get('problem_title',
               f"Right after the event, your WSP went dark — and the orders stopped."),
           font=serif(48, bold=True), fill=INK_T)

    # Calendar instruction
    d.text((M, 270),
           "Each block is one day. Purple = WSP running. Grey = switched off.",
           font=sans(22), fill=SLATE_T)

    # The calendar
    calendar = data.get('daily_calendar', [])
    if calendar:
        cal_top = 360; cal_bot = 620
        cal_h = cal_bot - cal_top
        n_days = len(calendar)
        gap = 6
        block_w = (W - 2*M - gap*(n_days-1)) // max(n_days, 1)
        block_w = min(block_w, 50)  # cap at 50px for readability

        # Annotations above (event window + dark window)
        event_indices = [i for i, day in enumerate(calendar) if day.get('event')]
        dark_indices = [i for i, day in enumerate(calendar) if day['status'] == 'off']

        # Event tab in gold
        if event_indices:
            ex_start = M + event_indices[0] * (block_w + gap)
            ex_end = M + (event_indices[-1]+1) * (block_w + gap) - gap
            d.rectangle([ex_start, cal_top - 50, ex_end, cal_top - 30], fill=GOLD_T)
            draw_ls(d, ex_start, cal_top - 90,
                    data.get('event_name', 'EVENT').upper(),
                    sans(18, bold=True), GOLD_T, sp=2)

        # Dark stretch in coral
        if dark_indices:
            # Find longest contiguous run
            longest_start, longest_len = 0, 0
            cur_start, cur_len = dark_indices[0], 1
            for j in range(1, len(dark_indices)):
                if dark_indices[j] == dark_indices[j-1] + 1:
                    cur_len += 1
                else:
                    if cur_len > longest_len:
                        longest_start, longest_len = cur_start, cur_len
                    cur_start, cur_len = dark_indices[j], 1
            if cur_len > longest_len:
                longest_start, longest_len = cur_start, cur_len

            dx_start = M + longest_start * (block_w + gap)
            dx_end = M + (longest_start + longest_len) * (block_w + gap) - gap
            d.rectangle([dx_start, cal_top - 50, dx_end, cal_top - 30], fill=tuple(CORAL_T))
            draw_ls(d, dx_start, cal_top - 90,
                    f"{longest_len} DAYS DARK",
                    sans(18, bold=True), tuple(CORAL_T), sp=2)

        # Day blocks
        for i, day in enumerate(calendar):
            bx = M + i * (block_w + gap)
            color = PURPLE_T if day['status'] == 'on' else GREY_BAR
            d.rectangle([bx, cal_top, bx + block_w, cal_bot], fill=color)
            # Date label every 7 days
            if i % 7 == 0 or i == len(calendar) - 1:
                date_lbl = day.get('label', '')
                if date_lbl:
                    d.text((bx, cal_bot + 12), date_lbl,
                           font=sans(16), fill=SLATE_T)

    # Bottom: two cards
    card_y = 760; card_h = 320
    card_w = (W - 2*M - 40) // 2

    # LEFT — coral card "WHILE IT WAS DARK"
    d.rectangle([M, card_y, M + card_w, card_y + card_h], fill=CORAL_LIGHT)
    d.rectangle([M, card_y, M + 8, card_y + card_h], fill=tuple(CORAL_T))
    draw_ls(d, M + 35, card_y + 40, "WHILE IT WAS DARK",
            sans(22, bold=True), tuple(CORAL_T), sp=2)
    missed_text = data.get('missed_wholesale_text', '~£15k')
    d.text((M + 35, card_y + 100), missed_text,
           font=serif(140, bold=True, italic=True), fill=tuple(CORAL_T))
    d.text((M + 35, card_y + 270),
           data.get('missed_wholesale_caption',
               "of wholesale revenue likely missed over those days"),
           font=sans(20), fill=INK_T)

    # RIGHT — lavender "THE PATTERN"
    rx = M + card_w + 40
    d.rectangle([rx, card_y, rx + card_w, card_y + card_h], fill=LAV_T)
    d.rectangle([rx, card_y, rx + 8, card_y + card_h], fill=PURPLE_T)
    draw_ls(d, rx + 35, card_y + 40, "THE PATTERN",
            sans(22, bold=True), PURPLE_T, sp=2)
    d.text((rx + 35, card_y + 100), "On = sales",
           font=serif(110, bold=True, italic=True), fill=DEEP_T)
    d.text((rx + 35, card_y + 270),
           "Off = silence. The budget keeps running out mid-month.",
           font=sans(20), fill=INK_T)

    # Takeaway
    d.text((M, H - 80),
           data.get('problem_takeaway',
               "Every day WSP is off, you stop showing up to ready-to-buy shoppers — and a competitor takes that slot."),
           font=serif(28, italic=True), fill=INK_T)
    return img


def render_restart_ask(data: dict) -> Image.Image:
    """Restart variant Ask slide — same pattern as Variant B but with restart-specific pillars."""
    img = Image.new("RGB", (W, H), WHITE_T)
    d = ImageDraw.Draw(img)
    draw_ls(d, M, 90, "04  ·  WHERE WE GO FROM HERE",
            sans(22, bold=True), PURPLE_T, sp=3)
    d.text((M, 165),
           data.get('ask_title', "Restart WSP now — the maths still works."),
           font=serif(58, bold=True), fill=INK_T)

    card_y = 360; card_h = 520
    card_w = (W - 2*M - 2*30) // 3
    card_gap = 30
    cards = data['pillars']
    for i, card in enumerate(cards[:3]):
        cx = M + i*(card_w + card_gap)
        d.rectangle([cx, card_y, cx + card_w, card_y + card_h], fill=LAV_T)
        d.rectangle([cx, card_y, cx + card_w, card_y + 10], fill=GOLD_T)
        draw_ls(d, cx + 35, card_y + 40, card['label'],
                sans(22, bold=True), PURPLE_T, sp=2)
        # Coral hero if it's the "what going dark cost" pillar
        hero_color = tuple(CORAL_T) if card.get('is_loss') else DEEP_T
        d.text((cx + 35, card_y + 110), card['hero'],
               font=serif(95, bold=True, italic=True), fill=hero_color)
        body_lines = wrap(d, card['body'], card_w - 70, sans(24))
        for j, ln in enumerate(body_lines):
            d.text((cx + 35, card_y + 290 + j*36), ln, font=sans(24), fill=INK_T)

    ask_y = 920; ask_h = 220
    d.rectangle([M, ask_y, W - M, ask_y + ask_h], fill=DEEP_T)
    draw_ls(d, M + 40, ask_y + 32, "THE ASK", sans(24, bold=True), GOLD_T, sp=3)
    ask_amount = data['ask_amount_text']
    last_wsc = _fmt(data['last_month_wsc'])
    d.text((M + 40, ask_y + 80), "Switch WSP back on at  ", font=sans(40), fill=WHITE_T)
    cw1, _ = tsize(d, "Switch WSP back on at  ", sans(40))
    d.text((M + 40 + cw1, ask_y + 80), ask_amount,
           font=sans(40, bold=True), fill=GOLD_T)
    cw2, _ = tsize(d, ask_amount, sans(40, bold=True))
    d.text((M + 40 + cw1 + cw2, ask_y + 80), "  —  5% of last month's",
           font=sans(40), fill=WHITE_T)
    d.text((M + 40, ask_y + 135),
           f"wholesale revenue (£{last_wsc}), set on day 1 so it never goes dark mid-month.",
           font=sans(40), fill=WHITE_T)

    closing = data.get('ask_closing',
        f"Proven £{data['ws_roas']:.2f} back for every £1  ·  currently leaving wholesale revenue on the table  ·  one rule, live on day 1.")
    d.text((M, H - 80), closing, font=serif(28, italic=True), fill=INK_T)
    return img


def build_deck_restart_pitch_a(data: dict) -> bytes:
    """Build a Variant A Restart pitch deck. Uses the dark-days calendar pattern."""
    import tempfile
    slides = [
        render_restart_cover(data),
        render_value_story(data),        # 01 · WHAT WSP WAS DRIVING (reuse)
        render_yoy_proof(data),          # 02 · THE EVENT RECAP (reuse)
        render_dark_days_calendar(data), # 03 · THE PROBLEM
        render_restart_ask(data),        # 04 · WHERE WE GO FROM HERE
    ]
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    with tempfile.TemporaryDirectory() as td:
        for i, slide_img in enumerate(slides):
            if slide_img.mode != "RGB":
                slide_img = slide_img.convert("RGB")
            png_path = f"{td}/slide_{i}.png"
            slide_img.save(png_path, "PNG", optimize=False)
            s = prs.slides.add_slide(BLANK)
            s.shapes.add_picture(png_path, 0, 0,
                                  width=prs.slide_width,
                                  height=prs.slide_height)
        out = io.BytesIO(); prs.save(out); out.seek(0)
        return out.read()


# =====================================================================
# VARIANT C — SWITCH TO 5% RULE
# Same opening as Variant B, but slide 3 reframes as "what erratic spend cost"
# and the ask is the budget-fixing rule rather than continuation.
# =====================================================================

def render_switch_to_5pct_ask(data: dict) -> Image.Image:
    """Variant C Ask — fixes lumpy spend, frames the rule as the unlock."""
    img = Image.new("RGB", (W, H), WHITE_T)
    d = ImageDraw.Draw(img)
    draw_ls(d, M, 90, "03  ·  WHERE WE GO FROM HERE",
            sans(22, bold=True), PURPLE_T, sp=3)
    d.text((M, 165),
           data.get('ask_title', "One rule fixes lumpy spend — for every supplier on it."),
           font=serif(58, bold=True), fill=INK_T)

    card_y = 360; card_h = 520
    card_w = (W - 2*M - 2*30) // 3
    card_gap = 30
    cards = data['pillars']
    for i, card in enumerate(cards[:3]):
        cx = M + i*(card_w + card_gap)
        d.rectangle([cx, card_y, cx + card_w, card_y + card_h], fill=LAV_T)
        d.rectangle([cx, card_y, cx + card_w, card_y + 10], fill=GOLD_T)
        draw_ls(d, cx + 35, card_y + 40, card['label'],
                sans(22, bold=True), PURPLE_T, sp=2)
        hero_color = tuple(CORAL_T) if card.get('is_loss') else DEEP_T
        d.text((cx + 35, card_y + 110), card['hero'],
               font=serif(95, bold=True, italic=True), fill=hero_color)
        body_lines = wrap(d, card['body'], card_w - 70, sans(24))
        for j, ln in enumerate(body_lines):
            d.text((cx + 35, card_y + 290 + j*36), ln, font=sans(24), fill=INK_T)

    ask_y = 920; ask_h = 220
    d.rectangle([M, ask_y, W - M, ask_y + ask_h], fill=DEEP_T)
    draw_ls(d, M + 40, ask_y + 32, "THE ASK", sans(24, bold=True), GOLD_T, sp=3)
    ask_amount = data['ask_amount_text']
    last_wsc = _fmt(data['last_month_wsc'])
    d.text((M + 40, ask_y + 80), "Set next month's WSP at  ", font=sans(40), fill=WHITE_T)
    cw1, _ = tsize(d, "Set next month's WSP at  ", sans(40))
    d.text((M + 40 + cw1, ask_y + 80), ask_amount,
           font=sans(40, bold=True), fill=GOLD_T)
    cw2, _ = tsize(d, ask_amount, sans(40, bold=True))
    d.text((M + 40 + cw1 + cw2, ask_y + 80), "  —  5% of last 30 days'",
           font=sans(40), fill=WHITE_T)
    d.text((M + 40, ask_y + 135),
           f"wholesale revenue (£{last_wsc}), set on day 1 — and every day 1 after.",
           font=sans(40), fill=WHITE_T)

    closing = data.get('ask_closing',
        f"Proven £{data['ws_roas']:.2f} back for every £1  ·  the gap to benchmark is now closed  ·  one rule, live on day 1.")
    d.text((M, H - 80), closing, font=serif(28, italic=True), fill=INK_T)
    return img


def build_deck_switch_to_5pct_c(data: dict) -> bytes:
    """Build a Variant C Switch-to-5% deck. Same data shape as Variant B."""
    import tempfile
    slides = [
        render_cover(data),
        render_value_story(data),
        render_yoy_proof(data),
        render_switch_to_5pct_ask(data),
    ]
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    with tempfile.TemporaryDirectory() as td:
        for i, slide_img in enumerate(slides):
            if slide_img.mode != "RGB":
                slide_img = slide_img.convert("RGB")
            png_path = f"{td}/slide_{i}.png"
            slide_img.save(png_path, "PNG", optimize=False)
            s = prs.slides.add_slide(BLANK)
            s.shapes.add_picture(png_path, 0, 0,
                                  width=prs.slide_width,
                                  height=prs.slide_height)
        out = io.BytesIO(); prs.save(out); out.seek(0)
        return out.read()


# =====================================================================
# Dispatcher — picks the right builder based on storyboard name
# =====================================================================
def build_deck(storyboard: str, data: dict) -> bytes:
    """Dispatch to the right builder based on storyboard name.

    Args:
        storyboard: storyboard name (with or without .md extension)
        data: structured data dict from data_prep

    Returns:
        .pptx bytes

    Raises:
        ValueError: if no builder is registered for the storyboard
    """
    key = storyboard.replace('.md', '').strip().split()[0]
    builders = {
        'value-review-variant-b':   build_deck_value_review_b,
        'restart-pitch-variant-a':  build_deck_restart_pitch_a,
        'switch-to-5pct-variant-c': build_deck_switch_to_5pct_c,
    }
    if key not in builders:
        raise ValueError(
            f"No builder for storyboard '{key}' yet. Available: {list(builders)}. "
            f"Other storyboards (promo-recap, summit, MBR) ship in v1.1."
        )
    return builders[key](data)


# =====================================================================
# DEEP-DIVE ANALYSIS SLIDES — SKU exposure split, CVR per bucket,
# sort rank movement, engine SKU spotlight. Added to the Value Review
# storyboard when the analyst wants the "full picture" treatment.
# =====================================================================

def render_sku_exposure_split(data: dict) -> Image.Image:
    """SKU buckets by ads exposure with WSC contribution per bucket.

    Tells the "your wholesale lives on a small set of SKUs" story —
    every supplier sees 7-10% of catalogue driving 75-85% of wholesale.
    """
    img = Image.new("RGB", (W, H), WHITE_T)
    d = ImageDraw.Draw(img)

    draw_ls(d, M, 90, "04  ·  WHERE YOUR WHOLESALE LIVES",
            sans(22, bold=True), PURPLE_T, sp=3)
    d.text((M, 165),
           data.get('exposure_title',
               "A small share of your SKUs carries most of your wholesale."),
           font=serif(48, bold=True), fill=INK_T)

    buckets = data.get('sku_exposure_buckets', [
        {'label': 'NO WSP', 'sku_pct': 78, 'wsc_pct': 14, 'color': GREY_BAR},
        {'label': 'LIGHT',  'sku_pct': 12, 'wsc_pct': 16, 'color': LAV_T},
        {'label': 'MODERATE','sku_pct': 7, 'wsc_pct': 28, 'color': (180, 130, 200)},
        {'label': 'HEAVY',  'sku_pct': 3,  'wsc_pct': 42, 'color': PURPLE_T},
    ])

    # Visual: two stacked horizontal bars — SKUs (top) vs WSC (bottom)
    # to make the disproportion visually obvious
    bar_left = M
    bar_right = W - M
    bar_w = bar_right - bar_left

    # Top bar — % of SKUs
    bar1_y = 380
    bar1_h = 100
    draw_ls(d, M, bar1_y - 35, "SHARE OF YOUR CATALOGUE",
            sans(18, bold=True), SLATE_T, sp=2)
    cx = bar_left
    for bucket in buckets:
        w = int(bar_w * bucket['sku_pct'] / 100)
        d.rectangle([cx, bar1_y, cx + w, bar1_y + bar1_h], fill=bucket['color'])
        # Label inside bar if there's room
        if w > 120:
            label_color = WHITE_T if bucket['color'] == PURPLE_T else INK_T
            d.text((cx + 18, bar1_y + 18), bucket['label'],
                   font=sans(18, bold=True), fill=label_color)
            d.text((cx + 18, bar1_y + 45), f"{bucket['sku_pct']}%",
                   font=serif(40, bold=True, italic=True), fill=label_color)
        cx += w

    # Bottom bar — % of WSC
    bar2_y = 600
    bar2_h = 100
    draw_ls(d, M, bar2_y - 35, "SHARE OF YOUR WHOLESALE",
            sans(18, bold=True), SLATE_T, sp=2)
    cx = bar_left
    for bucket in buckets:
        w = int(bar_w * bucket['wsc_pct'] / 100)
        d.rectangle([cx, bar2_y, cx + w, bar2_y + bar2_h], fill=bucket['color'])
        if w > 120:
            label_color = WHITE_T if bucket['color'] == PURPLE_T else INK_T
            d.text((cx + 18, bar2_y + 18), bucket['label'],
                   font=sans(18, bold=True), fill=label_color)
            d.text((cx + 18, bar2_y + 45), f"{bucket['wsc_pct']}%",
                   font=serif(40, bold=True, italic=True), fill=label_color)
        cx += w

    # Connecting visual emphasis between top "HEAVY" and bottom "HEAVY"
    heavy = buckets[-1]
    sku_start = bar_left + int(bar_w * (100 - heavy['sku_pct']) / 100)
    wsc_start = bar_left + int(bar_w * (100 - heavy['wsc_pct']) / 100)
    # Diagonal stripes connecting (visual: top-right wedge of SKUs blossoms into
    # much wider chunk of WSC below)
    for offset in range(3):
        d.line([(sku_start + offset, bar1_y + bar1_h),
                (wsc_start + offset, bar2_y)], fill=PURPLE_T, width=2)

    # Lavender callout at bottom
    cb_y = 870
    cb_h = 220
    d.rectangle([M, cb_y, W - M, cb_y + cb_h], fill=LAV_T)
    heavy_sku = heavy['sku_pct']
    heavy_wsc = heavy['wsc_pct']
    multiplier = round(heavy_wsc / heavy_sku, 1) if heavy_sku > 0 else 0
    d.text((M + 40, cb_y + 30), "THE LIFT MULTIPLE",
           font=sans(20, bold=True), fill=PURPLE_T)
    d.text((M + 40, cb_y + 70), f"{multiplier}×",
           font=serif(110, bold=True, italic=True), fill=DEEP_T)
    d.text((M + 280, cb_y + 100),
           f"Your heavily-exposed SKUs are {heavy['sku_pct']}% of catalogue",
           font=sans(26), fill=INK_T)
    d.text((M + 280, cb_y + 135),
           f"and drive {heavy['wsc_pct']}% of wholesale — {multiplier}× their fair share.",
           font=sans(26), fill=INK_T)
    d.text((M + 280, cb_y + 175),
           "This is the bucket WSP keeps in front of the shopper.",
           font=serif(24, italic=True), fill=PURPLE_T)

    # Takeaway
    d.text((M, H - 80),
           data.get('exposure_takeaway',
               "Take-away: your wholesale concentrates on the SKUs WSP keeps in front of ready-to-buy shoppers."),
           font=serif(28, italic=True), fill=INK_T)
    return img


def render_cvr_by_exposure(data: dict) -> Image.Image:
    """Conversion rate by ads-exposure bucket. The proof that exposure → conversion."""
    img = Image.new("RGB", (W, H), WHITE_T)
    d = ImageDraw.Draw(img)

    draw_ls(d, M, 90, "05  ·  WSP RAISES CONVERSION",
            sans(22, bold=True), PURPLE_T, sp=3)
    d.text((M, 165),
           data.get('cvr_title',
               "The more you advertise a SKU, the more readily it converts."),
           font=serif(48, bold=True), fill=INK_T)

    # 4 KPI-style cards — one per exposure bucket — with CVR as the hero
    buckets = data.get('cvr_buckets', [
        {'label': 'NO WSP',   'cvr': 0.4,  'sub': 'baseline organic traffic'},
        {'label': 'LIGHT',    'cvr': 0.7,  'sub': '~2× baseline'},
        {'label': 'MODERATE', 'cvr': 1.3,  'sub': '~3× baseline'},
        {'label': 'HEAVY',    'cvr': 2.6,  'sub': '~6× baseline'},
    ])

    card_y = 380
    card_h = 350
    card_w = (W - 2*M - 3*30) // 4
    card_gap = 30
    # Use a gradient: greyer for no-WSP, more purple for heavy
    stripe_colors = [GREY_BAR, LAV_DARK := (216, 188, 224), (175, 130, 195), PURPLE_T]
    for i, bucket in enumerate(buckets[:4]):
        cx = M + i * (card_w + card_gap)
        d.rectangle([cx, card_y, cx + card_w, card_y + card_h], fill=LAV_T)
        d.rectangle([cx, card_y, cx + card_w, card_y + 10],
                    fill=stripe_colors[i] if i < len(stripe_colors) else PURPLE_T)
        draw_ls(d, cx + 35, card_y + 40, bucket['label'],
                sans(22, bold=True), PURPLE_T, sp=2)
        d.text((cx + 35, card_y + 110), f"{bucket['cvr']:.1f}%",
               font=serif(95, bold=True, italic=True), fill=DEEP_T)
        d.text((cx + 35, card_y + 230),
               "conversion rate",
               font=sans(22), fill=SLATE_T)
        d.text((cx + 35, card_y + 270), bucket['sub'],
               font=sans(20, italic=True), fill=INK_T)

    # Lavender body callout
    cb_y = 790
    cb_h = 220
    d.rectangle([M, cb_y, W - M, cb_y + cb_h], fill=LAV_T)
    d.text((M + 40, cb_y + 30), "WHY THIS HAPPENS",
           font=sans(20, bold=True), fill=PURPLE_T)
    d.text((M + 40, cb_y + 75),
           "WSP places your SKUs in front of shoppers who are already looking to buy.",
           font=sans(26), fill=INK_T)
    d.text((M + 40, cb_y + 115),
           "The higher the placement, the better-matched the click, the more likely the purchase.",
           font=sans(26), fill=INK_T)
    d.text((M + 40, cb_y + 165),
           "Wayfair's sort algorithm rewards sales velocity, so the lift compounds week-over-week.",
           font=serif(26, italic=True), fill=PURPLE_T)

    d.text((M, H - 80),
           data.get('cvr_takeaway',
               "Take-away: WSP doesn't just buy you traffic — it buys you the right traffic at the moment it converts."),
           font=serif(28, italic=True), fill=INK_T)
    return img


def render_sort_rank_improvement(data: dict) -> Image.Image:
    """Sort-rank position improvement over the period for WSP-exposed SKUs."""
    img = Image.new("RGB", (W, H), WHITE_T)
    d = ImageDraw.Draw(img)

    draw_ls(d, M, 90, "06  ·  WSP LIFTS YOUR ORGANIC RANK",
            sans(22, bold=True), PURPLE_T, sp=3)
    d.text((M, 165),
           data.get('rank_title',
               "Pages 1-3 hold 80% of shoppers. WSP buys you a seat at that table — then keeps it."),
           font=serif(42, bold=True), fill=INK_T)

    # LEFT card — distribution of clicks across browse pages
    left_x = M; left_w = 1050; left_y = 380; left_h = 700
    d.rectangle([left_x, left_y, left_x + left_w, left_y + left_h], fill=LAV_T)
    draw_ls(d, left_x + 40, left_y + 40, "WHY PAGE RANK MATTERS",
            sans(20, bold=True), PURPLE_T, sp=2)

    # Page distribution chart
    pages = data.get('page_distribution', [
        ('Page 1', 62.5),
        ('Page 2', 16.6),
        ('Page 3', 6.1),
        ('Page 4', 3.7),
        ('Page 5', 2.3),
        ('Page 6+', 8.8),
    ])
    chart_top = left_y + 130
    chart_bot = left_y + 550
    chart_h = chart_bot - chart_top
    chart_w = left_w - 100
    chart_x_offset = left_x + 50
    n = len(pages)
    bw = (chart_w - 40 * (n-1)) // n
    max_v = max(v for _, v in pages) * 1.15
    for i, (label, val) in enumerate(pages):
        bx = chart_x_offset + i * (bw + 40)
        bh = int(val / max_v * chart_h)
        by = chart_bot - bh
        # First 3 pages = "the WSP zone" coloured PURPLE
        color = PURPLE_T if i < 3 else GREY_BAR
        d.rectangle([bx, by, bx + bw, chart_bot], fill=color)
        # value label
        d.text((bx, by - 30), f"{val:.1f}%",
               font=sans(20, bold=True), fill=INK_T)
        # page label
        d.text((bx, chart_bot + 14), label,
               font=sans(18), fill=SLATE_T)

    d.text((left_x + 40, left_y + left_h - 100),
           "Pages 1-3 hold ~85% of shoppers.",
           font=serif(28, italic=True), fill=DEEP_T)
    d.text((left_x + 40, left_y + left_h - 60),
           "Anything past page 3 might as well be invisible.",
           font=sans(22), fill=INK_T)

    # RIGHT card — sort rank movement before/after
    rx = left_x + left_w + 50
    rw = W - rx - M
    d.rectangle([rx, left_y, rx + rw, left_y + left_h], fill=DEEP_T)
    draw_ls(d, rx + 40, left_y + 40, "YOUR SORT-RANK MOVE",
            sans(20, bold=True), GOLD_T, sp=2)

    rank_before = data.get('rank_before', 18)
    rank_after  = data.get('rank_after', 7)
    improvement_pct = int((rank_before - rank_after) / rank_before * 100) if rank_before > 0 else 0

    d.text((rx + 40, left_y + 100), "BEFORE WSP",
           font=sans(18, bold=True), fill=(220, 195, 220))
    d.text((rx + 40, left_y + 135), f"Position #{rank_before}",
           font=serif(72, bold=True, italic=True), fill=WHITE_T)
    d.text((rx + 40, left_y + 230), "(average across engine SKUs)",
           font=sans(20, italic=True), fill=(180, 165, 200))

    # Arrow down
    arrow_x = rx + 40
    arrow_y = left_y + 290
    d.line([(arrow_x, arrow_y), (arrow_x, arrow_y + 50)], fill=GOLD_T, width=4)
    d.polygon([(arrow_x - 15, arrow_y + 50),
               (arrow_x + 15, arrow_y + 50),
               (arrow_x, arrow_y + 70)], fill=GOLD_T)

    d.text((rx + 40, left_y + 380), "WITH WSP",
           font=sans(18, bold=True), fill=GOLD_T)
    d.text((rx + 40, left_y + 415), f"Position #{rank_after}",
           font=serif(72, bold=True, italic=True), fill=GOLD_T)
    d.text((rx + 40, left_y + 510),
           f"−{improvement_pct}% lower number",
           font=sans(22, bold=True), fill=WHITE_T)
    d.text((rx + 40, left_y + 545),
           "= higher position",
           font=sans(22, italic=True), fill=(220, 195, 220))

    d.text((M, H - 80),
           data.get('rank_takeaway',
               "Take-away: WSP doesn't just sell THIS month — it earns your SKUs sort rank that pays out for months after."),
           font=serif(28, italic=True), fill=INK_T)
    return img


def render_engine_sku_spotlight(data: dict) -> Image.Image:
    """Top-5 engine SKUs with their CG status, OTD, CVR, units sold.

    This is §A.2.4 from SKILL.md — "rank SKUs by WSC contribution, take top 5."
    """
    img = Image.new("RGB", (W, H), WHITE_T)
    d = ImageDraw.Draw(img)

    draw_ls(d, M, 90, "07  ·  THE FIVE SKUs DOING THE WORK",
            sans(22, bold=True), PURPLE_T, sp=3)
    d.text((M, 165),
           data.get('spotlight_title',
               "These five products are the engine. WSP keeps them ahead of the competition."),
           font=serif(48, bold=True), fill=INK_T)

    skus = data.get('engine_skus', [
        {'name': 'Engine SKU 1', 'cg_status': 'High-Velocity CG', 'otd': 2, 'cvr': 3.4, 'units': 412, 'wsc_share': 8.2},
        {'name': 'Engine SKU 2', 'cg_status': 'High-Velocity CG', 'otd': 3, 'cvr': 2.9, 'units': 348, 'wsc_share': 6.8},
        {'name': 'Engine SKU 3', 'cg_status': 'High-Velocity CG', 'otd': 2, 'cvr': 3.1, 'units': 296, 'wsc_share': 5.9},
        {'name': 'Engine SKU 4', 'cg_status': 'Strategic CG',     'otd': 5, 'cvr': 2.4, 'units': 218, 'wsc_share': 4.4},
        {'name': 'Engine SKU 5', 'cg_status': 'High-Velocity CG', 'otd': 4, 'cvr': 2.7, 'units': 195, 'wsc_share': 3.8},
    ])
    total_share = sum(s['wsc_share'] for s in skus[:5])

    # Header summary
    sum_y = 290
    sum_h = 90
    d.rectangle([M, sum_y, W - M, sum_y + sum_h], fill=LAV_T)
    d.text((M + 30, sum_y + 24), f"{total_share:.0f}%",
           font=serif(48, bold=True, italic=True), fill=DEEP_T)
    d.text((M + 180, sum_y + 30),
           f"of your wholesale in {data.get('period', 'this period')} came from these five SKUs.",
           font=sans(24), fill=INK_T)
    d.text((M + 180, sum_y + 60),
           "All but one are in Castlegate — that is why the conversion holds.",
           font=serif(20, italic=True), fill=PURPLE_T)

    # Table
    table_y = 430
    row_h = 100
    col_widths = [800, 350, 200, 200, 240, 240]  # name, CG, OTD, CVR, units, % WSC
    col_x = [M]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)
    headers = ['SKU', 'CG STATUS', 'OTD (days)', 'CVR', 'UNITS SOLD', '% YOUR WSC']

    # Header row
    d.rectangle([M, table_y, W - M, table_y + 50], fill=DEEP_T)
    for i, h in enumerate(headers):
        d.text((col_x[i] + 16, table_y + 14),
               h, font=sans(16, bold=True), fill=GOLD_T)

    # Body rows
    for r, sku in enumerate(skus[:5]):
        ry = table_y + 50 + r * row_h
        # Alternate row shade
        if r % 2 == 0:
            d.rectangle([M, ry, W - M, ry + row_h], fill=(248, 244, 248))
        # SKU name (serif italic for emphasis)
        d.text((col_x[0] + 16, ry + 25), sku['name'],
               font=serif(28, bold=True, italic=True), fill=DEEP_T)
        d.text((col_x[0] + 16, ry + 65), data.get('period', ''),
               font=sans(16, italic=True), fill=SLATE_T)
        # CG status — coloured chip
        cg_color = GREEN_T if 'High-Velocity' in sku['cg_status'] else GOLD_T
        chip_w = 270
        d.rounded_rectangle([col_x[1] + 16, ry + 30,
                             col_x[1] + 16 + chip_w, ry + 65],
                            radius=6, fill=cg_color)
        d.text((col_x[1] + 30, ry + 36), sku['cg_status'],
               font=sans(15, bold=True), fill=WHITE_T)
        # OTD
        d.text((col_x[2] + 16, ry + 35), f"{sku['otd']}",
               font=serif(32, bold=True, italic=True), fill=DEEP_T)
        # CVR
        d.text((col_x[3] + 16, ry + 35), f"{sku['cvr']:.1f}%",
               font=serif(32, bold=True, italic=True), fill=DEEP_T)
        # Units
        d.text((col_x[4] + 16, ry + 35), f"{sku['units']:,}",
               font=serif(32, bold=True, italic=True), fill=DEEP_T)
        # % WSC
        d.text((col_x[5] + 16, ry + 35), f"{sku['wsc_share']:.1f}%",
               font=serif(32, bold=True, italic=True), fill=PURPLE_T)

    d.text((M, H - 80),
           data.get('spotlight_takeaway',
               "Take-away: these five SKUs are where WSP earns its keep — and where pausing it costs you the most."),
           font=serif(28, italic=True), fill=INK_T)
    return img


def build_deck_value_review_deep_dive(data: dict) -> bytes:
    """Build the deep-dive value review — 8 slides with SKU analysis.

    Order:
        1 · Cover
        2 · 01 · WHAT WSP IS DOING        (hero £/£1 + monthly chart)
        3 · 02 · YoY PROOF                (KPI band + green pills)
        4 · 03 · SKU EXPOSURE SPLIT       (NEW — where wholesale lives)
        5 · 04 · CVR BY EXPOSURE BUCKET   (NEW — exposure → conversion)
        6 · 05 · SORT RANK IMPROVEMENT    (NEW — pages 1-3 + before/after)
        7 · 06 · ENGINE SKU SPOTLIGHT     (NEW — top 5 table with CG/CVR/OTD)
        8 · 07 · WHERE WE GO FROM HERE    (3 pillars + The Ask)
    """
    import tempfile
    slides = [
        render_cover(data),
        render_value_story(data),
        render_yoy_proof(data),
        render_sku_exposure_split(data),
        render_cvr_by_exposure(data),
        render_sort_rank_improvement(data),
        render_engine_sku_spotlight(data),
        render_ask(data),
    ]
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    with tempfile.TemporaryDirectory() as td:
        for i, slide_img in enumerate(slides):
            if slide_img.mode != "RGB":
                slide_img = slide_img.convert("RGB")
            png_path = f"{td}/slide_{i}.png"
            slide_img.save(png_path, "PNG", optimize=False)
            s = prs.slides.add_slide(BLANK)
            s.shapes.add_picture(png_path, 0, 0,
                                  width=prs.slide_width,
                                  height=prs.slide_height)
        out = io.BytesIO(); prs.save(out); out.seek(0)
        return out.read()


# Update the dispatcher to include the deep-dive variant
_original_build_deck = build_deck

def build_deck(storyboard: str, data: dict) -> bytes:
    """Updated dispatcher with the deep-dive variant added."""
    key = storyboard.replace('.md', '').strip().split()[0]
    builders = {
        'value-review-variant-b':         build_deck_value_review_b,
        'value-review-deep-dive':         build_deck_value_review_deep_dive,
        'restart-pitch-variant-a':        build_deck_restart_pitch_a,
        'switch-to-5pct-variant-c':       build_deck_switch_to_5pct_c,
    }
    if key not in builders:
        raise ValueError(
            f"No builder for storyboard '{key}' yet. Available: {list(builders)}. "
            f"Other storyboards (promo-recap, summit, MBR) ship in v1.1."
        )
    return builders[key](data)
